"""Extract text content from uploaded files.

Supports: plain text, PDF, DOCX, PPTX.
Falls back gracefully for unsupported types.
"""

from __future__ import annotations

import asyncio
import io
import logging

from app.integrations.r2_store import R2Store

logger = logging.getLogger(__name__)

# Content types we can extract text from
_EXTRACTABLE_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Max chars to store (roughly ~8k tokens)
MAX_EXTRACTED_CHARS = 30_000


def can_extract(content_type: str) -> bool:
    """Check if we support text extraction for this content type."""
    return content_type in _EXTRACTABLE_TYPES


async def extract_text(
    store: R2Store,
    object_key: str,
    content_type: str,
) -> str | None:
    """Download a file from R2 and extract its text content.

    Returns None if extraction fails or the type is unsupported.
    """
    if not can_extract(content_type):
        return None

    try:
        data = await asyncio.to_thread(store.download_object, object_key)
    except Exception:
        logger.exception("Failed to download %s for text extraction", object_key)
        return None

    try:
        text = _extract_from_bytes(data, content_type)
    except Exception:
        logger.exception("Failed to extract text from %s (type=%s)", object_key, content_type)
        return None

    if not text or not text.strip():
        return None

    text = text.strip()
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS] + "\n\n[... tekst afgekapt ...]"

    return text


def _extract_from_bytes(data: bytes, content_type: str) -> str | None:
    if content_type in ("text/plain", "text/markdown", "text/csv"):
        return _extract_plain_text(data)
    if content_type == "application/pdf":
        return _extract_pdf(data)
    if content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx(data)
    if content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_pptx(data)
    return None


def _extract_plain_text(data: bytes) -> str:
    for encoding in ("utf-8", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    import pymupdf

    pages: list[str] = []
    with pymupdf.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text = page.get_text()
            if isinstance(text, str) and text.strip():
                pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)
